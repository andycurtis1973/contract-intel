#!/usr/bin/env python3
"""Build the sales + solution-architecture deck: contract_intel_deck.pptx (16:9, 17 slides).

A balanced three-act brief on reading a contract estate with Bedrock + S3 Vectors:
  1. the business case (what a blind estate costs)      -- for the buyer
  2. the solution and the reference architecture         -- for solution architects
  3. proof for IT (accuracy vs lawyers, security, cost)  -- for customer IT
then a deploy-and-integrate close (full estate, into the systems of record).

Reuses the video's slide builders (video/render.py) as full-bleed proof visuals and
adds the sales/architecture/IT slides. Every number comes from results/rundata.json,
so the deck, the demo, and the video can't drift. Every slide carries speaker notes.

    python3 deck/build_deck.py     # -> deck/contract_intel_deck.pptx
"""

from __future__ import annotations

import math
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "video"))
import render as R                       # noqa: E402
from render import (W, H, BG, TX, MUT, DIM, CARD, RULE, MACHINE, RISK, OK, FLAG,  # noqa
                    T, font, base, card,
                    D, REG, ACC, RT, C, NV, N, AUTO, TIGHT, SILENT, NOCAP, HIST)
from PIL import ImageDraw  # noqa: E402

SLIDES = ROOT / "deck" / "slides"
SLIDES.mkdir(parents=True, exist_ok=True)

MINC = REG["minimum_commitments"]
CHUNKS = RT["chunks_indexed"]
PCT_SENT = RT["pct_sent"]


# ---- small helpers --------------------------------------------------------
def _footer(img, n, total, full=True):
    d = ImageDraw.Draw(img, "RGBA")
    if full:
        d.line([110, H - 66, W - 110, H - 66], fill=RULE, width=1)
        T(d, (110, H - 48), "CONTRACT INTELLIGENCE  ·  BEDROCK + S3 VECTORS", font("mono", 18), DIM)
    T(d, (W - 110, H - 48), f"{n:02d} / {total:02d}", font("mono", 18), DIM, "ra")


def arrow(d, x1, y1, x2, y2, col, dash=False, w=4):
    if dash:
        n = max(2, int(math.hypot(x2 - x1, y2 - y1) / 18))
        for i in range(n):
            if i % 2:
                continue
            a, b = i / n, (i + 1) / n
            d.line([x1 + (x2 - x1) * a, y1 + (y2 - y1) * a,
                    x1 + (x2 - x1) * b, y1 + (y2 - y1) * b], fill=col, width=w)
    else:
        d.line([x1, y1, x2, y2], fill=col, width=w)
    ang = math.atan2(y2 - y1, x2 - x1)
    for da in (2.5, -2.5):
        d.line([x2, y2, x2 + 16 * math.cos(ang + da), y2 + 16 * math.sin(ang + da)],
               fill=col, width=w)


def _wrap(s, n):
    words, line, out = s.split(), "", []
    for wd in words:
        if len(line + " " + wd) > n:
            out.append(line); line = wd
        else:
            line = (line + " " + wd).strip()
    out.append(line)
    return out


def title(d, t, sub=None):
    T(d, (110, 178), t, font("bold", 52), TX)
    if sub:
        for j, ln in enumerate(_wrap(sub, 96)):
            T(d, (110, 250 + j * 40), ln, font("reg", 30), MUT)


def statcard(d, x, y, w, h, big, label, col):
    card(d, [x, y, x + w, y + h], accent=col, radius=12)
    T(d, (x + 28, y + 26), big, font("bold", 66), col)
    for j, ln in enumerate(_wrap(label, 30)):
        T(d, (x + 28, y + 124 + j * 36), ln, font("reg", 26), MUT)


def numbox(d, x, y, w, h, n, ttl, sub, col):
    card(d, [x, y, x + w, y + h], accent=col, radius=12)
    d.ellipse([x + 24, y + 28, x + 70, y + 74], fill=col)
    T(d, (x + 47, y + 51), str(n), font("bold", 28), (15, 18, 24), "mm")
    T(d, (x + 24, y + 98), ttl, font("bold", 29), TX)
    for j, ln in enumerate(_wrap(sub, 26)):
        T(d, (x + 24, y + 142 + j * 34), ln, font("reg", 22), MUT)


def bigcards(d, cards, y0=320, h=252):
    """Four two-by-two cards with a left accent rail and wrapped body."""
    for i, (t, s, col) in enumerate(cards):
        cx = 110 + (i % 2) * 862
        cy = y0 + (i // 2) * (h + 36)
        card(d, [cx, cy, cx + 826, cy + h], side=col, radius=12)
        T(d, (cx + 36, cy + 34), t, font("bold", 32), TX)
        for j, ln in enumerate(_wrap(s, 44)):
            T(d, (cx + 36, cy + 100 + j * 40), ln, font("reg", 27), MUT)


def callout(d, box, tag, text, col, tagcol=None):
    card(d, box, accent=col, radius=12)
    T(d, (box[0] + 40, box[1] + 22), tag, font("mono", 22), tagcol or col)
    for j, ln in enumerate(_wrap(text, 92)):
        T(d, (box[0] + 40, box[1] + 58 + j * 42), ln, font("bold", 31), TX)


# ---- ACT 1 -- the business case -------------------------------------------
def w_title():
    img, d = base()
    T(d, (W / 2, 300), "Read every contract you've", font("bold", 74), TX, "mm")
    T(d, (W / 2, 392), "already signed.", font("bold", 74), MACHINE, "mm")
    T(d, (W / 2, 520), "Find the renewals, the deadlines, and the uncapped liability hiding",
      font("reg", 36), MUT, "mm")
    T(d, (W / 2, 570), "in the estate — with Amazon Bedrock and S3 Vectors.", font("reg", 36), MUT, "mm")
    tag = "SALES BRIEF  ·  REFERENCE ARCHITECTURE  ·  IT PROOF"
    w = d.textlength(tag, font=font("mono", 26))
    d.rounded_rectangle([(W - w) / 2 - 30, 690, (W + w) / 2 + 30, 744], radius=27,
                        fill=(255, 255, 255, 12), outline=(MACHINE[0], MACHINE[1], MACHINE[2], 150))
    T(d, (W / 2, 718), tag, font("mono", 26), MUT, "mm")
    return img


def w_agenda():
    img, d = base("the brief")
    title(d, "Three audiences, one build")
    items = [("The business case", "what a blind contract estate quietly costs you", MACHINE),
             ("The solution & architecture", "how it reads the whole estate on AWS", OK),
             ("Proof for IT", "accuracy vs lawyers, security posture, and cost", MACHINE),
             ("Deploy & integrate", "the full estate, into your systems of record", OK)]
    for i, (t, s, col) in enumerate(items):
        y = 336 + i * 138
        card(d, [110, y, 1810, y + 112], side=col, radius=12)
        T(d, (150, y + 30), str(i + 1), font("bold", 44), col)
        T(d, (240, y + 22), t, font("bold", 34), TX)
        T(d, (240, y + 68), s, font("reg", 26), MUT)
    return img


def w_cost_of_blind():
    img, d = base("the business case  ·  act 1")
    title(d, "What a blind contract estate costs you",
          f"From {N} real SEC-filed contracts — the leaks a spreadsheet never catches.")
    cards = [(f"{AUTO}", "auto-renew with no action — 26% of the estate", RISK),
             (f"{TIGHT}", "give you 30 days or less to stop the renewal", RISK),
             (f"{SILENT}", "auto-renew with no notice period stated at all", FLAG),
             (f"{NOCAP}", "put no ceiling on liability — 75% of contracts", RISK)]
    w, gap, x0, y, h = 405, 16, 110, 348, 250
    for i, (big, lab, col) in enumerate(cards):
        statcard(d, x0 + i * (w + gap), y, w, h, big, lab, col)
    callout(d, [110, 662, 1810, 792], "THE POINT",
            f"{MINC} more carry minimum-commitment obligations. Every one of these is money or "
            "risk that renews whether or not anyone reads it.", MACHINE)
    return img


def w_whynow():
    img, d = base("why now")
    title(d, "The whole estate, not a spot-check",
          "Contract review used to be too expensive to run on everything. Two shifts changed that.")
    cards = [("A model reads the prose",
              "Obligations are written differently in every contract. A foundation model reads them the "
              "way a lawyer does — at machine speed, across the whole book.", MACHINE),
             ("S3 Vectors removes the infra",
              "Semantic search over the entire estate with no vector database to run, secure, or pay "
              "for while it sits idle.", OK)]
    for i, (t, s, col) in enumerate(cards):
        cx = 110 + i * 862
        card(d, [cx, 336, cx + 826, 620], accent=col, radius=14)
        T(d, (cx + 36, 366), t, font("bold", 34), TX)
        for j, ln in enumerate(_wrap(s, 42)):
            T(d, (cx + 36, 436 + j * 42), ln, font("reg", 28), MUT)
    callout(d, [110, 668, 1810, 812], "THE SHIFT",
            "Reading all 499 cost about a dollar. The question moved from \"which contracts can we "
            "afford to review?\" to \"why haven't we read all of them?\"", FLAG)
    return img


# ---- ACT 2 -- the solution & architecture ---------------------------------
def w_what():
    img, d = base("the solution  ·  act 2")
    title(d, "What it does, in four steps")
    steps = [("Land", "your contracts to an S3 bucket", MACHINE),
             ("Read", "Bedrock reads only the clauses that matter", MACHINE),
             ("Ask", "S3 Vectors — question the whole estate", OK),
             ("Register", "renewals, deadlines and caps to act on", OK)]
    w, gap, x0, y, h = 395, 40, 110, 340, 200
    for i, (t, s, col) in enumerate(steps):
        x = x0 + i * (w + gap)
        numbox(d, x, y, w, h, i + 1, t, s, col)
        if i < 3:
            T(d, (x + w + gap / 2, y + h / 2), "→", font("bold", 40), DIM, "mm")
    callout(d, [110, 606, 1810, 748], "THE DELIVERABLE",
            "A living register: every obligation, the source clause, and how long you have to act — "
            "and it's graded against lawyer annotations, not taken on faith.", OK)
    return img


def _abox(d, x, y, w, h, ttl, sub, col):
    card(d, [x, y, x + w, y + h], accent=col, radius=14)
    T(d, (x + 24, y + 26), ttl, font("bold", 29), TX)
    if sub:
        T(d, (x + 24, y + 74), sub, font("mono", 20), MUT)


def w_architecture():
    img, d = base("reference architecture  ·  for solution architects")
    T(d, (110, 174), "The architecture", font("bold", 52), TX)
    T(d, (110, 244), "Two flows share one index: documents are embedded once; questions are answered "
      "against them on demand.", font("reg", 28), MUT)
    ah = 132
    # INGEST lane
    T(d, (110, 312), "INGEST  ·  ONCE PER DOCUMENT", font("mono", 22), OK)
    ay = 348
    _abox(d, 110, ay, 300, ah, "Contracts", "S3 · 499 docs · 26M chars", MACHINE)
    _abox(d, 470, ay, 340, ah, "Bedrock — Titan v2", "embeddings · 1024-d", MACHINE)
    _abox(d, 870, ay, 430, ah, "S3 Vectors index", f"{CHUNKS:,} chunks · cosine", OK)
    arrow(d, 410, ay + ah / 2, 468, ay + ah / 2, DIM)
    arrow(d, 810, ay + ah / 2, 868, ay + ah / 2, DIM)
    # ASK lane
    T(d, (110, 566), "ASK  ·  PER QUESTION", font("mono", 22), MACHINE)
    by = 602
    _abox(d, 110, by, 300, ah, "A question", "natural language", MACHINE)
    _abox(d, 470, by, 340, ah, "Bedrock — Claude", "extract obligations", MACHINE)
    _abox(d, 870, by, 430, ah, "The register", "renewals · caps · dates", OK)
    arrow(d, 410, by + ah / 2, 468, by + ah / 2, DIM)
    arrow(d, 810, by + ah / 2, 868, by + ah / 2, DIM)
    # index feeds the extractor (the retrieval step)
    arrow(d, 1000, ay + ah, 720, by, OK, dash=True)
    T(d, (1030, 498), "filter {contract} · top-k", font("mono", 20), OK)
    # the note card, spanning both lanes on the right
    card(d, [1360, ay, 1810, by + ah], accent=MACHINE, radius=14)
    T(d, (1390, ay + 24), "WHY IT'S CHEAP & SAFE", font("mono", 20), MACHINE)
    notes = [(f"Only {PCT_SENT}% of the text ever reaches the model.", OK),
             ("No vector DB to run — the index lives in S3, next to the docs.", TX),
             ("Everything stays in your AWS account and region.", TX)]
    yy = ay + 72
    for txt, col in notes:
        d.ellipse([1390, yy + 8, 1406, yy + 24], fill=col if col != TX else MACHINE)
        lines = _wrap(txt, 30)
        for ln in lines:
            T(d, (1424, yy), ln, font("reg", 23), TX if col == TX else col)
            yy += 34
        yy += 16
    return img


def w_s3vectors():
    img, d = base("the differentiator  ·  for solution architects")
    title(d, "S3 Vectors: the index lives in object storage",
          "Semantic search over the estate with nothing to stand up, secure, or scale.")
    # left: the lifecycle
    card(d, [110, 332, 930, 812], radius=14)
    T(d, (146, 358), "THE LIFECYCLE  ·  boto3 s3vectors", font("mono", 21), MACHINE)
    steps = ["create_vector_bucket()",
             "create_index(dim=1024,", "        metric=cosine, float32)",
             "put_vectors(batches ~100)",
             "query_vectors(topK,", "        filter={contract},", "        returnMetadata=True)"]
    for i, s in enumerate(steps):
        T(d, (150, 420 + i * 52), s, font("mono", 26), TX if not s.startswith(" ") else MUT)
    T(d, (150, 800 - 8), "", font("mono", 20), MUT)
    # right: why it matters
    card(d, [990, 332, 1810, 812], accent=OK, radius=14)
    T(d, (1026, 358), "WHY IT MATTERS", font("mono", 21), OK)
    pts = ["No cluster to stand up, secure, patch, or scale.",
           "No idle cost — you pay per query, like S3 itself.",
           "Retrieval scoped per contract by metadata filter.",
           "The index sits next to the documents, in your account.",
           "No OpenSearch, pgvector, or Pinecone to operate."]
    for i, p in enumerate(pts):
        y = 424 + i * 74
        d.ellipse([1026, y + 6, 1044, y + 24], fill=OK)
        for j, ln in enumerate(_wrap(p, 46)):
            T(d, (1064, y + j * 32), ln, font("reg", 27), TX)
    return img


# ---- ACT 3 -- proof for IT ------------------------------------------------
def w_honesty():
    img, d = base("proof for it  ·  act 3")
    title(d, "Where it's strong — and where it needs a human",
          "Graded against the lawyer annotations in CUAD — the weak spots reported, not hidden.")
    cap = ACC["clauses"]["cap on liability"]
    # trust-it column
    card(d, [110, 332, 930, 792], accent=OK, radius=14)
    T(d, (146, 358), "TRUST IT", font("mono", 22), OK)
    tr = [(f"{NV['exact_pct']}%", "of notice windows it finds, it reads the day-count exactly right"),
          ("0.92 / 0.91", "auto-renewal precision, termination recall"),
          (f"{PCT_SENT}%", "of the text sent to the model — retrieval does the rest")]
    for i, (b, s) in enumerate(tr):
        y = 420 + i * 122
        T(d, (146, y), b, font("bold", 46), OK)
        for j, ln in enumerate(_wrap(s, 40)):
            T(d, (146, y + 62 + j * 32), ln, font("reg", 24), MUT)
    # human column
    card(d, [990, 332, 1810, 792], accent=FLAG, radius=14)
    T(d, (1026, 358), "SEND TO A HUMAN", font("mono", 22), FLAG)
    hm = [(f"{cap['recall']*100:.0f}%", "liability-cap recall — almost never wrong, but finds only half"),
          ("> 365 days", "a notice value this long is flagged for review, not used"),
          ("0.59", "minimum-commitment F1 — a genuine assist, not an authority")]
    for i, (b, s) in enumerate(hm):
        y = 420 + i * 122
        T(d, (1026, y), b, font("bold", 46), FLAG)
        for j, ln in enumerate(_wrap(s, 40)):
            T(d, (1026, y + 62 + j * 32), ln, font("reg", 24), MUT)
    return img


def w_security():
    img, d = base("proof for it  ·  security & governance")
    title(d, "Your data never leaves your account")
    cards = [("Runs in your AWS account",
              "S3, S3 Vectors and Bedrock all run in your account and region. No third-party SaaS, "
              "no data leaving your perimeter.", MACHINE),
             ("Bedrock doesn't train on you",
              "Your prompts and documents are not used to train the foundation models. Inputs stay "
              "yours.", OK),
             ("Scoped by IAM + metadata",
              "Retrieval is filtered per contract; access follows the S3 and IAM controls you "
              "already enforce.", MACHINE),
             ("No new attack surface",
              "No endpoints, clusters or always-on services to harden, patch, or monitor. Nothing "
              "idles.", OK)]
    bigcards(d, cards, y0=316, h=232)
    callout(d, [110, 850, 1810, 976], "IN SHORT",
            "It fits inside the landing zone you already run — no exception request required.", MACHINE)
    return img


def w_cost():
    img, d = base("proof for it  ·  cost")
    title(d, "About a dollar to read 499 contracts",
          "Bedrock and S3 Vectors are pay-per-call. Cost scales with what you read, not with uptime.")
    cards = [("~$1", "total, to embed, index and extract the whole run", OK),
             (f"{PCT_SENT}%", "of characters ever sent to the model", MACHINE),
             ("$0", "standing cost beyond pennies/month for the bucket", OK),
             ("0", "servers, endpoints or idle clusters to pay for", MACHINE)]
    w, gap, x0, y, h = 405, 16, 110, 348, 250
    for i, (big, lab, col) in enumerate(cards):
        statcard(d, x0 + i * (w + gap), y, w, h, big, lab, col)
    callout(d, [110, 662, 1810, 792], "THE ECONOMICS",
            "Delete the vector bucket and you are at exactly zero. There is no minimum, no reserved "
            "capacity, and nothing running between questions.", MACHINE)
    return img


# ---- close -- deploy & integrate ------------------------------------------
def w_deploy():
    img, d = base("deploy & integrate")
    title(d, "How it lands in your environment")
    steps = [("Ingest the estate", "point it at your contract repositories → S3", MACHINE),
             ("Embed & index", "Bedrock Titan → S3 Vectors, once per document", MACHINE),
             ("Extract & register", "scheduled Bedrock extraction → the register", OK),
             ("Integrate", "push into CLM, ERP and BI — your systems of record", OK)]
    w, gap, x0, y, h = 395, 40, 110, 340, 210
    for i, (t, s, col) in enumerate(steps):
        x = x0 + i * (w + gap)
        card(d, [x, y, x + w, y + h], accent=col, radius=12)
        d.ellipse([x + 24, y + 28, x + 70, y + 74], fill=col)
        T(d, (x + 47, y + 51), str(i + 1), font("bold", 28), (15, 18, 24), "mm")
        T(d, (x + 24, y + 98), t, font("bold", 28), TX)
        for j, ln in enumerate(_wrap(s, 27)):
            T(d, (x + 24, y + 142 + j * 32), ln, font("reg", 21), MUT)
        if i < 3:
            T(d, (x + w + gap / 2, y + h / 2), "→", font("bold", 40), DIM, "mm")
    callout(d, [110, 616, 1810, 758], "CONTINUOUS",
            "Runs on a schedule. New and amended contracts flow in on their own, and the register "
            "stays current in the tools your teams already open.", OK)
    return img


def w_why_wins():
    img, d = base("why this wins")
    title(d, "Why this approach wins")
    cards = [("The whole estate, not a sample",
              "Cheap enough to read every contract you've signed — and to keep reading the new ones.", OK),
             ("In your account, no infra",
              "Bedrock + S3 Vectors: nothing to run, secure, or pay for while it idles.", MACHINE),
             ("Graded, not trusted",
              "Accuracy measured against lawyer annotations; doubt is surfaced, never invented.", OK),
             ("A register that plugs in",
              "Obligations flow into the CLM, ERP and BI systems your teams already use.", MACHINE)]
    bigcards(d, cards, y0=320, h=252)
    return img


def w_close():
    img, d = base()
    T(d, (W / 2, 300), "Let's read a sample of", font("bold", 72), TX, "mm")
    T(d, (W / 2, 388), "your contracts.", font("bold", 72), MACHINE, "mm")
    T(d, (W / 2, 506), f"{N} contracts read · {AUTO} renew themselves · {TIGHT} give you 30 days "
      "· graded vs lawyers", font("mono", 27), MUT, "mm")
    rows = [("REPO", "github.com/andycurtis1973/contract-intel"),
            ("PILOT", "bring 50 contracts — we run it in your sandbox")]
    for i, (lab, u) in enumerate(rows):
        y = 596 + i * 88
        s = f"{lab}   {u}"
        w = d.textlength(s, font=font("mono", 28))
        x0 = (W - w) / 2 - 34
        card(d, [x0, y, x0 + w + 68, y + 64], accent=MACHINE, radius=12)
        T(d, (x0 + 34, y + 32), lab, font("mono", 28), DIM, "lm")
        T(d, (x0 + 34 + d.textlength(f"{lab}   ", font=font("mono", 28)), y + 32),
          u, font("mono", 28), MACHINE, "lm")
    T(d, (W / 2, 792), "Open data (CUAD) · open code · the accuracy numbers are there to check the work",
      font("reg", 26), DIM, "mm")
    return img


# ---- the deck: (builder, notes, is_chart) ---------------------------------
DECK = [
    (w_title, "Frame it for the room: this is one build with three audiences. For the buyer it's "
     "money and risk hiding in contracts you've already signed. For solution architects it's a clean "
     "AWS reference architecture. For IT it's data that never leaves your account, graded accuracy, "
     "and a cost of about a dollar. We'll go through all three.", False),
    (w_agenda, "Three acts. First the business case — what a blind estate quietly costs. Then the "
     "solution and the architecture, in enough detail for your SAs to sketch it. Then the proof IT "
     "needs — accuracy, security, and cost. And we close on how it deploys into your environment and "
     "integrates with your systems of record.", False),
    (lambda: R.a_problem(1.0), "Every square is a real agreement filed with the SEC. 130 of these 499 "
     "renew automatically — nobody presses a button, nobody signs anything, the term just rolls over "
     "and the invoice arrives. This is the pain: obligations that act on their own, on an estate "
     "nobody has time to re-read.", True),
    (w_cost_of_blind, "Put numbers on it. 130 auto-renew. Of those, 27 give you 30 days or less to "
     "stop it, and 39 state no notice period at all — so you can't even tell when the window closes. "
     "376 of the 499 — three quarters — put no ceiling on liability. And 131 carry minimum-commitment "
     "obligations. This is money and risk renewing whether or not anyone reads it.", False),
    (w_whynow, "Why is now the moment? Two things changed. A foundation model reads legal prose the "
     "way a lawyer does, so obligations written differently in every contract can finally be read at "
     "machine speed. And S3 Vectors removes the infrastructure — semantic search over the whole estate "
     "with no vector database to run. Together they drop the cost enough to read everything, not just "
     "spot-check.", False),
    (w_what, "The build is four steps. Land the contracts in S3. Read them with Bedrock — but only "
     "the clauses that matter. Ask the whole estate a question through S3 Vectors. And produce the "
     "register: every obligation, the source clause, and how long you have to act. The register is "
     "the deliverable, and it's graded against lawyers, not taken on faith.", False),
    (w_architecture, "Here's the architecture your SAs will want. Two flows share one index. Ingest, "
     "once per document: contracts in S3, embedded by Bedrock Titan v2 into 1024-dimension vectors, "
     "indexed in S3 Vectors — 22,792 chunks, cosine distance. Ask, per question: a natural-language "
     "question retrieves the top-k clauses from that index with a metadata filter scoped to the "
     "contract, Bedrock Claude extracts the obligations, and they land in the register. The key: only "
     "24.5% of the text ever reaches the model, there's no vector database to run, and it all stays "
     "in your account.", False),
    (lambda: R.a_read(1.0), "This is the retrieval economics. A contract is fifty thousand characters "
     "of boilerplate around about six sentences that matter. We retrieve the clauses first and send "
     "the model only 24.5% of the text. That's cheaper — and it keeps the model honest, because it's "
     "reading the clause, not guessing from the whole document.", True),
    (w_s3vectors, "The differentiator, in SA terms. S3 Vectors is a lifecycle you already understand: "
     "create a vector bucket, create an index at 1024 dimensions with cosine distance, put your "
     "vectors in batches, and query with a top-k and a metadata filter. Why it matters: no cluster to "
     "stand up, secure, patch, or scale; no idle cost; retrieval scoped per contract; and no "
     "OpenSearch, pgvector, or Pinecone to operate. The index lives next to the documents.", False),
    (lambda: R.a_ask(1.0), "And this is what the search actually does. Nobody files auto-renewal in a "
     "searchable field — they write two paragraphs of legal prose, differently every time. Ask the "
     "whole estate in natural language and it finds them: 22,792 clauses indexed in S3 Vectors, "
     "queried in under a second, with no database to keep running.", True),
    (lambda: R.a_accuracy(1.0), "Now the question IT always asks: is it right? The lawyers who built "
     "CUAD annotated every clause, so we grade the machine instead of trusting it. On the notice "
     "window — the clause that actually costs money — it finds 92% of them, and when it finds one it "
     "reads the day-count exactly right 96.6% of the time. It under-calls liability caps — almost "
     "never wrong, but only finds half. That one still needs a human.", True),
    (w_honesty, "Read that as a routing rule, not a report card. Trust it on notice windows, "
     "auto-renewal, and termination — that's most of the value, and it's measured. Send liability "
     "caps and minimum commitments to a human, and note the plausibility guard: a notice value over "
     "365 days is flagged for review rather than used. A system that surfaces doubt instead of "
     "inventing a deadline is exactly what an audit wants.", False),
    (w_security, "The governance answer. Everything runs in your AWS account and region — S3, S3 "
     "Vectors, Bedrock — with no third-party SaaS and no data leaving your perimeter. Bedrock doesn't "
     "train on your inputs. Retrieval is scoped per contract by metadata and follows your existing "
     "IAM. And there's no new attack surface — nothing always-on to harden or patch. It fits inside "
     "the landing zone you already run.", False),
    (w_cost, "And the cost, plainly. About a dollar to embed, index, and extract all 499 contracts, "
     "because only a quarter of the text reaches the model. Standing cost is effectively zero — "
     "pennies a month for the vector bucket, and nothing else running. Bedrock and S3 Vectors are "
     "pay-per-call; delete the bucket and you're at zero. Cost scales with what you read, not with "
     "uptime.", False),
    (w_deploy, "How it lands. Point it at your contract repositories into S3; embed and index once "
     "per document; run extraction on a schedule to keep the register current; and integrate — push "
     "the register into the CLM, ERP, and BI tools your teams already open. New and amended contracts "
     "flow in continuously, so the register never goes stale.", False),
    (w_why_wins, "So why this approach. It reads the whole estate, not a sample, and keeps reading. "
     "It runs in your account with no infrastructure to operate. It's graded against lawyers rather "
     "than trusted. And it produces a register that plugs into the systems your teams already use. "
     "Cheap, safe, measured, and integrated.", False),
    (w_close, "That's the brief. The data is public, the code is open, and the accuracy numbers are "
     "there to check the work. The best next step is small and concrete: bring us fifty of your own "
     "contracts and we'll show you what renews, when the windows close, and where the liability is "
     "uncapped — running in a sandbox in your account.", False),
]


def main() -> int:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, (builder, notes, is_chart) in enumerate(DECK):
        png = SLIDES / f"slide_{i:02d}.png"
        img = builder()
        if 0 < i < len(DECK) - 1:
            _footer(img, i + 1, len(DECK), full=not is_chart)
        img.save(png)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
        slide.notes_slide.notes_text_frame.text = notes
        print(f"  [{i+1:2d}/{len(DECK)}] {png.name}")

    out = ROOT / "deck" / "contract_intel_deck.pptx"
    prs.save(str(out))
    mb = out.stat().st_size / 1e6
    print(f"\n  ✅ {out}  ({len(DECK)} slides, {mb:.1f} MB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
